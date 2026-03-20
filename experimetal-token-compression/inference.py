from __future__ import annotations

import argparse
import math
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, List, Optional

from llmlingua import PromptCompressor
from docx import Document
from pptx import Presentation
from pypdf import PdfReader


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt"}


@dataclass
class TextUnit:
    source_name: str
    unit_id: str
    text: str


def normalize_text(text: str) -> str:
    if not text:
        return ""

    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = "\n".join(line.rstrip() for line in text.splitlines())
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]{2,}", " ", text)
    return text.strip()


def extract_text_from_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def extract_text_from_docx(path: Path) -> str:
    doc = Document(str(path))
    parts: List[str] = []

    for p in doc.paragraphs:
        text = p.text.strip()
        if text:
            parts.append(text)

    for table in doc.tables:
        for row in table.rows:
            row_cells = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                if cell_text:
                    row_cells.append(cell_text)
            if row_cells:
                parts.append(" | ".join(row_cells))

    return "\n".join(parts)


def iter_shape_text(shape) -> Iterable[str]:
    if getattr(shape, "has_text_frame", False):
        text = shape.text.strip()
        if text:
            yield text

    if getattr(shape, "has_table", False):
        table = shape.table
        for row in table.rows:
            row_cells = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                if cell_text:
                    row_cells.append(cell_text)
            if row_cells:
                yield " | ".join(row_cells)

    if hasattr(shape, "shapes"):
        for subshape in shape.shapes:
            yield from iter_shape_text(subshape)


def extract_text_from_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    parts: List[str] = []

    for idx, slide in enumerate(prs.slides, start=1):
        slide_parts: List[str] = [f"[Slide {idx}]"]

        for shape in slide.shapes:
            for text in iter_shape_text(shape):
                slide_parts.append(text)

        try:
            notes_slide = slide.notes_slide
            notes = []
            for shape in notes_slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = shape.text.strip()
                    if text:
                        notes.append(text)
            if notes:
                slide_parts.append("[Notes]")
                slide_parts.extend(notes)
        except Exception:
            pass

        parts.append("\n".join(slide_parts))

    return "\n\n".join(parts)


def extract_text_from_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    parts: List[str] = []

    for page_num, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""

        text = text.strip()
        if text:
            parts.append(f"[Page {page_num}]\n{text}")

    return "\n\n".join(parts)


def extract_text_from_file(path_str: str) -> str:
    path = Path(path_str).expanduser().resolve()
    if not path.exists():
        raise FileNotFoundError(f"Input file does not exist: {path}")

    ext = path.suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(
            f"Unsupported file type: {ext}. Supported: {sorted(SUPPORTED_EXTENSIONS)}"
        )

    if ext == ".txt":
        raw = extract_text_from_txt(path)
    elif ext == ".docx":
        raw = extract_text_from_docx(path)
    elif ext == ".pptx":
        raw = extract_text_from_pptx(path)
    elif ext == ".pdf":
        raw = extract_text_from_pdf(path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

    cleaned = normalize_text(raw)
    if not cleaned:
        raise ValueError(f"No extractable text found in {path.name}")
    return cleaned


def split_text_into_chunks(text: str, max_chars: int = 4000, overlap_chars: int = 300) -> List[str]:
    """
    Chunk long text to avoid feeding huge single contexts.
    Uses paragraph-aware splitting first, then falls back to hard slicing.
    """
    if len(text) <= max_chars:
        return [text]

    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
    chunks: List[str] = []
    current = ""

    for para in paragraphs:
        if not current:
            current = para
            continue

        candidate = current + "\n\n" + para
        if len(candidate) <= max_chars:
            current = candidate
        else:
            chunks.append(current)
            if overlap_chars > 0 and len(current) > overlap_chars:
                tail = current[-overlap_chars:]
                current = tail + "\n\n" + para
            else:
                current = para

    if current:
        chunks.append(current)

    final_chunks: List[str] = []
    for chunk in chunks:
        if len(chunk) <= max_chars:
            final_chunks.append(chunk)
        else:
            start = 0
            step = max_chars - overlap_chars if max_chars > overlap_chars else max_chars
            while start < len(chunk):
                final_chunks.append(chunk[start:start + max_chars].strip())
                start += step

    return [c for c in final_chunks if c.strip()]


def load_inputs(
    text: Optional[str],
    files: Optional[List[str]],
    chunk_chars: int,
    overlap_chars: int,
) -> List[TextUnit]:
    units: List[TextUnit] = []

    if text and files:
        raise ValueError("Provide either --text or --files, not both.")
    if not text and not files:
        raise ValueError("Provide one input source: --text or --files.")

    if text:
        cleaned = normalize_text(text)
        chunks = split_text_into_chunks(cleaned, max_chars=chunk_chars, overlap_chars=overlap_chars)
        for i, chunk in enumerate(chunks, start=1):
            units.append(
                TextUnit(
                    source_name="raw_text",
                    unit_id=f"raw_text::chunk_{i}",
                    text=chunk,
                )
            )
        return units

    assert files is not None
    for file_path in files:
        path = Path(file_path).expanduser().resolve()
        extracted = extract_text_from_file(str(path))
        chunks = split_text_into_chunks(extracted, max_chars=chunk_chars, overlap_chars=overlap_chars)
        for i, chunk in enumerate(chunks, start=1):
            units.append(
                TextUnit(
                    source_name=path.name,
                    unit_id=f"{path.name}::chunk_{i}",
                    text=chunk,
                )
            )
    return units


def choose_batch_size(num_input_files: int, num_units: int) -> int:
    """
    Variable batch sizing:
    - fewer files => larger batches
    - more files => smaller batches
    This is a simple heuristic to preserve boundaries and keep memory stable.
    """
    if num_units <= 1:
        return 1
    if num_input_files <= 2:
        return min(8, num_units)
    if num_input_files <= 5:
        return min(6, num_units)
    if num_input_files <= 10:
        return min(4, num_units)
    return min(2, num_units)


def batched(items: List[TextUnit], batch_size: int) -> List[List[TextUnit]]:
    return [items[i:i + batch_size] for i in range(0, len(items), batch_size)]


def build_compressor(model_name: str, device_map: str) -> PromptCompressor:
    return PromptCompressor(
        model_name=model_name,
        use_llmlingua2=True,
        device_map=device_map,
    )


def compress_batch(
    compressor: PromptCompressor,
    batch_units: List[TextUnit],
    rate: float,
    target_token: int,
) -> dict:
    contexts = [u.text for u in batch_units]

    results = compressor.compress_prompt(
        contexts,
        rate=rate,
        target_token=target_token,
        use_context_level_filter=False,
        use_token_level_filter=True,
        keep_split=True,
        force_tokens=["\n"],
    )
    return results


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Batch-aware LLMLingua-2 compression for raw text, PDF, DOCX, PPTX, and TXT."
    )
    parser.add_argument("--text", type=str, default=None, help="Raw input text.")
    parser.add_argument(
        "--files",
        nargs="+",
        default=None,
        help="One or more files: .txt .docx .pptx .pdf",
    )
    parser.add_argument(
        "--device",
        type=str,
        default="cpu",
        help="Device map: cpu, cuda, mps, auto, etc.",
    )
    parser.add_argument(
        "--rate",
        type=float,
        default=0.33,
        help="Compression rate target.",
    )
    parser.add_argument(
        "--target-token",
        type=int,
        default=-1,
        help="Optional explicit token target per batch.",
    )
    parser.add_argument(
        "--chunk-chars",
        type=int,
        default=4000,
        help="Max characters per chunk before batching.",
    )
    parser.add_argument(
        "--overlap-chars",
        type=int,
        default=300,
        help="Overlap between adjacent chunks.",
    )
    parser.add_argument(
        "--save-dir",
        type=str,
        default=None,
        help="Optional directory to save extracted and compressed chunks.",
    )

    args = parser.parse_args()

    num_input_files = 1 if args.text else len(args.files)
    units = load_inputs(
        text=args.text,
        files=args.files,
        chunk_chars=args.chunk_chars,
        overlap_chars=args.overlap_chars,
    )

    if not units:
        raise ValueError("No text units were produced.")

    batch_size = choose_batch_size(num_input_files=num_input_files, num_units=len(units))
    batches = batched(units, batch_size=batch_size)

    compressor = build_compressor(
        model_name="microsoft/llmlingua-2-bert-base-multilingual-cased-meetingbank",
        device_map=args.device,
    )

    save_dir = Path(args.save_dir).resolve() if args.save_dir else None
    if save_dir:
        save_dir.mkdir(parents=True, exist_ok=True)

    total_origin_tokens = 0
    total_compressed_tokens = 0
    all_compressed_outputs: List[str] = []

    print(f"Total text units: {len(units)}")
    print(f"Adaptive batch size: {batch_size}")
    print(f"Total batches: {len(batches)}")
    print("-" * 80)

    for batch_idx, batch_units in enumerate(batches, start=1):
        results = compress_batch(
            compressor=compressor,
            batch_units=batch_units,
            rate=args.rate,
            target_token=args.target_token,
        )

        compressed_prompt = results["compressed_prompt"]
        origin_tokens = results["origin_tokens"]
        compressed_tokens = results["compressed_tokens"]
        compression_rate = results["rate"]

        total_origin_tokens += origin_tokens
        total_compressed_tokens += compressed_tokens

        unit_labels = [u.unit_id for u in batch_units]

        print(f"Batch {batch_idx}/{len(batches)}")
        print(f"Sources: {unit_labels}")
        print(f"Compressed prompt: {compressed_prompt}")
        print(f"Original tokens: {origin_tokens}")
        print(f"Compressed tokens: {compressed_tokens}")
        print(f"Compression rate: {compression_rate}")
        print("-" * 80)

        all_compressed_outputs.append(compressed_prompt)

        if save_dir:
            batch_file = save_dir / f"batch_{batch_idx:03d}_compressed.txt"
            meta_file = save_dir / f"batch_{batch_idx:03d}_sources.txt"
            batch_file.write_text(compressed_prompt, encoding="utf-8")
            meta_file.write_text("\n".join(unit_labels), encoding="utf-8")

    aggregate_rate = (
        total_compressed_tokens / total_origin_tokens if total_origin_tokens > 0 else 0.0
    )

    print("FINAL AGGREGATE")
    print(f"Compressed prompt: {'\n\n'.join(all_compressed_outputs)}")
    print(f"Original tokens: {total_origin_tokens}")
    print(f"Compressed tokens: {total_compressed_tokens}")
    print(f"Compression rate: {aggregate_rate}")


if __name__ == "__main__":
    main()