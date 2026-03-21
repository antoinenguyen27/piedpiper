from __future__ import annotations

import re
from dataclasses import dataclass
from io import BytesIO

from .runtime import get_prompt_compressor
from .schemas import CompressionItemResult, TextOptions

TEXT_SUFFIXES = {".txt", ".md", ".pdf", ".docx", ".pptx"}


@dataclass(slots=True)
class SourceText:
    item_id: str
    index: int
    source_name: str
    text: str


@dataclass(slots=True)
class TextUnit:
    item_id: str
    index: int
    source_name: str
    unit_id: str
    text: str


def normalize_text(text: str) -> str:
    if not text:
        return ""

    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = "\n".join(line.rstrip() for line in text.splitlines())
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]{2,}", " ", text)
    return text.strip()


def extract_text_from_txt_bytes(data: bytes) -> str:
    return data.decode("utf-8", errors="ignore")


def extract_text_from_docx_bytes(data: bytes) -> str:
    from docx import Document

    doc = Document(BytesIO(data))
    parts: list[str] = []

    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)

    for table in doc.tables:
        for row in table.rows:
            row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_cells:
                parts.append(" | ".join(row_cells))

    return "\n".join(parts)


def iter_shape_text(shape):
    if getattr(shape, "has_text_frame", False):
        text = shape.text.strip()
        if text:
            yield text

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_cells:
                yield " | ".join(row_cells)

    if hasattr(shape, "shapes"):
        for subshape in shape.shapes:
            yield from iter_shape_text(subshape)


def extract_text_from_pptx_bytes(data: bytes) -> str:
    from pptx import Presentation

    presentation = Presentation(BytesIO(data))
    parts: list[str] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_parts = [f"[Slide {slide_index}]"]
        for shape in slide.shapes:
            slide_parts.extend(iter_shape_text(shape))

        try:
            notes_slide = slide.notes_slide
        except Exception:  # pragma: no cover - library-dependent branch
            notes_slide = None

        if notes_slide is not None:
            notes = [
                shape.text.strip()
                for shape in notes_slide.shapes
                if getattr(shape, "has_text_frame", False) and shape.text.strip()
            ]
            if notes:
                slide_parts.append("[Notes]")
                slide_parts.extend(notes)

        parts.append("\n".join(slide_parts))

    return "\n\n".join(parts)


def extract_text_from_pdf_bytes(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(BytesIO(data))
    parts: list[str] = []

    for page_number, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""

        text = text.strip()
        if text:
            parts.append(f"[Page {page_number}]\n{text}")

    return "\n\n".join(parts)


def extract_text(source_name: str, data: bytes) -> str:
    lower_name = source_name.lower()
    if lower_name.endswith((".txt", ".md")):
        raw = extract_text_from_txt_bytes(data)
    elif lower_name.endswith(".docx"):
        raw = extract_text_from_docx_bytes(data)
    elif lower_name.endswith(".pptx"):
        raw = extract_text_from_pptx_bytes(data)
    elif lower_name.endswith(".pdf"):
        raw = extract_text_from_pdf_bytes(data)
    else:
        raise ValueError(f"Unsupported text file type for extraction: {source_name}")

    cleaned = normalize_text(raw)
    if not cleaned:
        raise ValueError(f"No extractable text found in {source_name}")
    return cleaned


def split_text_into_chunks(
    text: str,
    *,
    max_chars: int = 4000,
    overlap_chars: int = 300,
) -> list[str]:
    if len(text) <= max_chars:
        return [text]

    paragraphs = [paragraph.strip() for paragraph in text.split("\n\n") if paragraph.strip()]
    chunks: list[str] = []
    current = ""

    for paragraph in paragraphs:
        if not current:
            current = paragraph
            continue

        candidate = f"{current}\n\n{paragraph}"
        if len(candidate) <= max_chars:
            current = candidate
            continue

        chunks.append(current)
        if overlap_chars > 0 and len(current) > overlap_chars:
            current = f"{current[-overlap_chars:]}\n\n{paragraph}"
        else:
            current = paragraph

    if current:
        chunks.append(current)

    final_chunks: list[str] = []
    for chunk in chunks:
        if len(chunk) <= max_chars:
            final_chunks.append(chunk)
            continue

        step = max_chars - overlap_chars if max_chars > overlap_chars else max_chars
        start = 0
        while start < len(chunk):
            final_chunks.append(chunk[start : start + max_chars].strip())
            start += step

    return [chunk for chunk in final_chunks if chunk]


def choose_batch_size(*, num_input_files: int, num_units: int) -> int:
    if num_units <= 1:
        return 1
    if num_input_files <= 2:
        return min(8, num_units)
    if num_input_files <= 5:
        return min(6, num_units)
    if num_input_files <= 10:
        return min(4, num_units)
    return min(2, num_units)


def batched(items: list[TextUnit], batch_size: int) -> list[list[TextUnit]]:
    return [items[i : i + batch_size] for i in range(0, len(items), batch_size)]


def build_units(sources: list[SourceText], options: TextOptions) -> list[TextUnit]:
    units: list[TextUnit] = []
    for source in sources:
        chunks = split_text_into_chunks(
            source.text,
            max_chars=options.chunk_chars,
            overlap_chars=options.overlap_chars,
        )
        for chunk_index, chunk in enumerate(chunks, start=1):
            units.append(
                TextUnit(
                    item_id=source.item_id,
                    index=source.index,
                    source_name=source.source_name,
                    unit_id=f"{source.item_id}::chunk_{chunk_index}",
                    text=chunk,
                )
            )
    return units


def compress_text_sources(
    sources: list[SourceText],
    options: TextOptions,
) -> tuple[list[CompressionItemResult], dict[str, int]]:
    if not sources:
        return [], {"origin_tokens": 0, "compressed_tokens": 0}

    units = build_units(sources, options)
    batch_size = choose_batch_size(num_input_files=len(sources), num_units=len(units))
    batches = batched(units, batch_size=batch_size)
    compressor = get_prompt_compressor()

    results_by_item: dict[str, dict[str, object]] = {
        source.item_id: {
            "index": source.index,
            "source_name": source.source_name,
            "chunks": [],
            "origin_tokens": 0,
            "compressed_tokens": 0,
        }
        for source in sources
    }
    usage = {"origin_tokens": 0, "compressed_tokens": 0}

    for batch in batches:
        payload = compressor.compress_prompt(
            [unit.text for unit in batch],
            rate=options.rate,
            target_token=options.target_token,
            use_context_level_filter=False,
            use_token_level_filter=True,
            keep_split=True,
            force_tokens=["\n"],
        )

        compressed_contexts = payload.get("compressed_prompt_list")
        if not compressed_contexts:
            compressed_contexts = [payload["compressed_prompt"]]

        for unit, compressed_text in zip(batch, compressed_contexts):
            item = results_by_item[unit.item_id]
            item["chunks"].append((unit.unit_id, compressed_text))

        origin_tokens = int(payload.get("origin_tokens", 0))
        compressed_tokens = int(payload.get("compressed_tokens", 0))
        usage["origin_tokens"] += origin_tokens
        usage["compressed_tokens"] += compressed_tokens

        if batch:
            origin_share = max(1, origin_tokens // len(batch))
            compressed_share = max(1, compressed_tokens // len(batch))
            for unit in batch:
                item = results_by_item[unit.item_id]
                item["origin_tokens"] += origin_share
                item["compressed_tokens"] += compressed_share

    items: list[CompressionItemResult] = []
    for source in sorted(sources, key=lambda entry: entry.index):
        item = results_by_item[source.item_id]
        ordered_chunks = [chunk_text for _, chunk_text in item["chunks"]]
        origin_tokens = int(item["origin_tokens"])
        compressed_tokens = int(item["compressed_tokens"])
        rate = (
            compressed_tokens / origin_tokens
            if origin_tokens > 0
            else 0.0
        )
        items.append(
            CompressionItemResult(
                id=source.item_id,
                index=source.index,
                modality="text",
                source_name=source.source_name,
                status="completed",
                output_text="\n\n".join(ordered_chunks).strip(),
                metrics={
                    "origin_tokens": origin_tokens,
                    "compressed_tokens": compressed_tokens,
                    "compression_rate": round(rate, 4),
                    "chunk_count": len(ordered_chunks),
                    "batch_size": batch_size,
                },
            )
        )

    return items, usage

